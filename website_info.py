import os
from pptx import Presentation
from datetime import datetime
from pptx.util import Inches, Pt

today = datetime.today().strftime("%Y%m%d")

prs = Presentation()

#P1 title page
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]

f = open('site_title.txt')
title.text = f.read()
f = open('site_subtitle.txt')
subtitle.text = f.read()

#P2 Website product lists
bullet_slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(bullet_slide_layout)
shapes = slide.shapes

title_shape = shapes.title
body_shape = shapes.placeholders[1]

title_shape.text = 'Your Website information'

tf = body_shape.text_frame
tf.text = 'Your Web site products list'

#p = tf.add_paragraph()
f = open('product_lists.txt')
if len(f.read()) is None:
    p = tf.add_paragraph()
    p.text = 'There is no infomation.'
    p.level = 2
else:
    f = open('product_lists.txt')
    products = list(map(str.strip,(f.read().split("\n"))))
    for product in products:
        p = tf.add_paragraph()
        p.text = product
        p.level = 2
        p.font.size = Pt(17)

#P3 MX record
bullet_slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(bullet_slide_layout)
shapes = slide.shapes

title_shape = shapes.title
body_shape = shapes.placeholders[1]

title_shape.text = 'Your Web site information'

tf = body_shape.text_frame
tf.text = 'Your Web site MX record'

#p = tf.add_paragraph()
f = open('MX_record.txt')
if len(f.read()) == 0:
    p = tf.add_paragraph()
    p.text = 'There is no infomation.'
    p.level = 2
else:
    f = open('MX_record.txt')
    products = list(map(str.strip,(f.read().split("\n"))))
    for product in products:
        p = tf.add_paragraph()
        p.text = product
        p.level = 1
        p.font.size = Pt(17)

#P4 WHOIS record
bullet_slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(bullet_slide_layout)
shapes = slide.shapes

title_shape = shapes.title
body_shape = shapes.placeholders[1]

title_shape.text = 'Your Web site information'

tf = body_shape.text_frame
tf.text = 'Your Web site WHOIS record'

#p = tf.add_paragraph()
f = open('WHOIS_record.txt')
if len(f.read()) == 0:
    p = tf.add_paragraph()
    p.text = 'There is no infomation.'
    p.level = 1
else:
    f = open('WHOIS_record.txt')
    p = tf.add_paragraph()
    p.text = f.read()
    p.level = 1
    p.font.size = Pt(10)


#P5 site top page
img_path = 'image.jpg'

blank_slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(blank_slide_layout)

left = width = Inches(2.5)
top = Inches(2)
height = Inches(3.5)
pic = slide.shapes.add_picture(img_path, left, top, height=height)

left = Inches(2.3)
top = width = height = Inches(5.2)
txBox = slide.shapes.add_textbox(left, top, width, height)
tf = txBox.text_frame

tf.text = ""

p = tf.add_paragraph()
p.text = "your site TOP Page"
p.font.size = Pt(40)

prs.save('website_info.pptx')
